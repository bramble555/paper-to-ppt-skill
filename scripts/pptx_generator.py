from __future__ import annotations

import re
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _hex_to_rgb(hex_color: str, default: tuple[int, int, int]) -> tuple[int, int, int]:
    m = re.match(r"#?([0-9a-fA-F]{6})", hex_color or "")
    if not m:
        return default
    v = m.group(1)
    return int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16)


def generate_pptx_from_html(html_path: Path, output_path: Path) -> None:
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "html.parser")
    sections = soup.find_all("section", class_="slide")

    prs = Presentation()
    for sec in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        style = sec.get("style", "")
        bg = "#141b2d"
        if "background:" in style:
            bg = style.split("background:", 1)[1].split(";", 1)[0].strip()
        r, g, b = _hex_to_rgb(bg, (20, 27, 45))
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

        title_tag = sec.find(["h1", "h2"])
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = title_tag.get_text(strip=True) if title_tag else "Untitled"
        tp.font.name = "Calibri"
        tp.font.size = Pt(34)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(232, 236, 255)
        tp.alignment = PP_ALIGN.LEFT

        ul = sec.find("ul")
        body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.8))
        body = body_box.text_frame
        body.clear()
        if ul:
            for idx, li in enumerate(ul.find_all("li")):
                p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
                p.text = li.get_text(" ", strip=True)
                p.level = 0
                p.font.name = "Calibri"
                p.font.size = Pt(22)
                p.font.color.rgb = RGBColor(232, 236, 255)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
